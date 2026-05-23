"""Build the NOMIN AI study platform pitch deck following the
Customer Court Battle - Step 1 Group Presentations template."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# Brand palette inspired by the UE template (mint background, deep red text)
MINT = RGBColor(0xBF, 0xE2, 0xDF)
DEEP_RED = RGBColor(0x8C, 0x1F, 0x28)
INK = RGBColor(0x1F, 0x1F, 0x1F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xF5, 0xF1, 0xEC)
ACCENT = RGBColor(0x2B, 0x6B, 0x6B)


def set_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if not line:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = DEEP_RED
        shape.line.width = Pt(1.25)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, *, size=16, color=INK, bullet="•"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_footer(slide, page_num, total):
    # Left side branding
    add_text(slide, Inches(0.4), Inches(7.05), Inches(5), Inches(0.3),
             "NOMIN  ·  AI Study Buddy", size=10, bold=True, color=DEEP_RED)
    # Right side page number
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.8), Inches(0.3),
             f"{page_num} / {total}", size=10, color=DEEP_RED, align=PP_ALIGN.RIGHT)


def add_side_bar(slide):
    # Vertical accent bar on left, like the template's red stripe
    add_rect(slide, Emu(0), Emu(0), Inches(0.25), Inches(7.5), DEEP_RED)


def make_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, MINT)
    add_side_bar(slide)

    # Header tag
    add_text(slide, Inches(0.7), Inches(0.6), Inches(8), Inches(0.4),
             "CUSTOMER COURT BATTLE  ·  STEP 1 GROUP PRESENTATION",
             size=14, bold=True, color=DEEP_RED)

    # Big title
    add_text(slide, Inches(0.7), Inches(1.8), Inches(12), Inches(1.6),
             "NOMIN", size=96, bold=True, color=DEEP_RED)
    add_text(slide, Inches(0.7), Inches(3.5), Inches(12), Inches(0.8),
             "Your AI Study Buddy", size=36, bold=True, color=INK)
    add_text(slide, Inches(0.7), Inches(4.3), Inches(12), Inches(0.6),
             "Personalized learning that actually fits how you study.",
             size=20, color=INK)

    # Team strip
    add_rect(slide, Inches(0.7), Inches(5.6), Inches(12), Inches(0.05), DEEP_RED)
    add_text(slide, Inches(0.7), Inches(5.75), Inches(8), Inches(0.4),
             "University of Europe for Applied Sciences  ·  Entrepreneurship Course",
             size=12, bold=True, color=DEEP_RED)
    add_text(slide, Inches(0.7), Inches(6.15), Inches(8), Inches(0.4),
             "Presented by Team NOMIN", size=12, color=INK)

    add_footer(slide, 1, 8)
    return slide


def make_section_slide(prs, number, kicker, title, body_renderer, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, MINT)
    add_side_bar(slide)

    # Step badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(0.55),
                                   Inches(0.9), Inches(0.9))
    badge.fill.solid()
    badge.fill.fore_color.rgb = DEEP_RED
    badge.line.fill.background()
    btf = badge.text_frame
    btf.margin_left = Emu(0); btf.margin_right = Emu(0)
    btf.margin_top = Emu(0); btf.margin_bottom = Emu(0)
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = btf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = str(number)
    run.font.size = Pt(32); run.font.bold = True; run.font.color.rgb = WHITE

    add_text(slide, Inches(1.8), Inches(0.55), Inches(11), Inches(0.45),
             kicker, size=12, bold=True, color=DEEP_RED)
    add_text(slide, Inches(1.8), Inches(0.95), Inches(11), Inches(0.7),
             title, size=30, bold=True, color=INK)

    # Divider
    add_rect(slide, Inches(0.7), Inches(1.85), Inches(12), Inches(0.04), DEEP_RED)

    body_renderer(slide)

    add_footer(slide, page, 8)
    return slide


def slide_idea(slide):
    # Big statement card
    card = add_rect(slide, Inches(0.7), Inches(2.2), Inches(12), Inches(1.8), WHITE)
    add_text(slide, Inches(1.0), Inches(2.45), Inches(11.4), Inches(1.4),
             "NOMIN is an AI study buddy that turns lecture notes, slides,\n"
             "and textbooks into a personal tutor — available 24/7, in any\n"
             "language, on any device.",
             size=22, bold=True, color=INK)

    # Three pillars
    pillar_y = Inches(4.25)
    pillar_h = Inches(2.4)
    gap = Inches(0.2)
    pw = Inches((12 - 0.4) / 3 - 0.2)
    titles = ["Understand", "Practice", "Master"]
    descs = [
        "Upload any material — NOMIN summarizes it in your style and "
        "explains it the way you learn best.",
        "Auto-generates quizzes, flashcards, and mock exams from your "
        "course content.",
        "Tracks your weak spots, schedules spaced repetition, and gets "
        "you exam-ready faster."
    ]
    for i, (t, d) in enumerate(zip(titles, descs)):
        x = Inches(0.7) + i * (pw + gap)
        add_rect(slide, x, pillar_y, pw, pillar_h, WHITE)
        add_rect(slide, x, pillar_y, pw, Inches(0.4), DEEP_RED)
        add_text(slide, x + Inches(0.2), pillar_y + Inches(0.05),
                 pw - Inches(0.4), Inches(0.35), t,
                 size=14, bold=True, color=WHITE)
        add_text(slide, x + Inches(0.2), pillar_y + Inches(0.55),
                 pw - Inches(0.4), Inches(1.7), d, size=13, color=INK)


def slide_customer(slide):
    # Hero statement
    add_text(slide, Inches(0.7), Inches(2.15), Inches(12), Inches(0.5),
             "Not everyone. Our primary customer is:",
             size=16, color=INK)
    add_text(slide, Inches(0.7), Inches(2.55), Inches(12), Inches(0.7),
             "University students aged 18–26 in Europe",
             size=28, bold=True, color=DEEP_RED)
    add_text(slide, Inches(0.7), Inches(3.3), Inches(12), Inches(0.5),
             "studying STEM, business, or design programs — laptop-first, "
             "phone-second, exam-stressed.",
             size=14, color=INK)

    # 5 personas
    personas = [
        ("Lena, 21", "Bachelor — Business",
         "Juggles part-time job. Skims slides the night before."),
        ("Arjun, 23", "Master — Computer Science",
         "Drowns in dense papers. Wants a tutor that speaks his level."),
        ("Sofía, 19", "Foundation Year",
         "ESL student. Needs concepts re-explained in her language."),
        ("Marek, 25", "Online MBA",
         "Full-time job. Studies in 20-minute commute chunks."),
        ("Yuki, 22", "Bachelor — Design",
         "Visual learner. Hates walls of text, loves diagrams."),
    ]
    pw = Inches(2.42)
    gap = Inches(0.06)
    py = Inches(4.3)
    ph = Inches(2.4)
    for i, (name, role, desc) in enumerate(personas):
        x = Inches(0.7) + i * (pw + gap)
        add_rect(slide, x, py, pw, ph, WHITE)
        add_rect(slide, x, py, pw, Inches(0.45), DEEP_RED)
        add_text(slide, x + Inches(0.15), py + Inches(0.05),
                 pw - Inches(0.3), Inches(0.4), name,
                 size=13, bold=True, color=WHITE)
        add_text(slide, x + Inches(0.15), py + Inches(0.55),
                 pw - Inches(0.3), Inches(0.4), role,
                 size=11, bold=True, color=DEEP_RED)
        add_text(slide, x + Inches(0.15), py + Inches(0.95),
                 pw - Inches(0.3), Inches(1.4), desc,
                 size=11, color=INK)


def slide_need(slide):
    add_text(slide, Inches(0.7), Inches(2.15), Inches(12), Inches(0.5),
             "Their main need is…", size=16, color=INK)
    add_text(slide, Inches(0.7), Inches(2.55), Inches(12), Inches(1.5),
             "…to actually understand their course material —\n"
             "and pass exams — without burning out.",
             size=28, bold=True, color=DEEP_RED)

    # Problem cards
    problems = [
        ("Information overload",
         "Hundreds of slides per course, no clear path through them."),
        ("Generic content",
         "YouTube tutors and textbooks aren't tailored to their syllabus."),
        ("Expensive tutoring",
         "Private tutors cost €30–80/hour — out of reach for most students."),
        ("Procrastination loop",
         "Without structure, study time slides into TikTok time."),
    ]
    cw = Inches(5.9)
    ch = Inches(1.2)
    gap = Inches(0.2)
    base_x = Inches(0.7)
    base_y = Inches(4.4)
    for i, (title, desc) in enumerate(problems):
        col = i % 2
        row = i // 2
        x = base_x + col * (cw + gap)
        y = base_y + row * (ch + gap)
        add_rect(slide, x, y, cw, ch, WHITE)
        add_rect(slide, x, y, Inches(0.15), ch, DEEP_RED)
        add_text(slide, x + Inches(0.35), y + Inches(0.12),
                 cw - Inches(0.5), Inches(0.4), title,
                 size=14, bold=True, color=DEEP_RED)
        add_text(slide, x + Inches(0.35), y + Inches(0.55),
                 cw - Inches(0.5), Inches(0.6), desc,
                 size=12, color=INK)


def slide_impact(slide):
    add_text(slide, Inches(0.7), Inches(2.15), Inches(12), Inches(0.5),
             "By helping this customer, our start-up also…", size=16, color=INK)

    # Headline equation
    add_text(slide, Inches(0.7), Inches(2.6), Inches(12), Inches(0.9),
             "Reduces drop-out rates  by  giving every student a 1:1 tutor.",
             size=24, bold=True, color=DEEP_RED)
    add_text(slide, Inches(0.7), Inches(3.4), Inches(12), Inches(0.5),
             "Improves grades  by  closing knowledge gaps before exams.",
             size=24, bold=True, color=DEEP_RED)

    # Stats strip
    stats = [
        ("32%", "of EU students drop out before graduating"),
        ("€2.4B", "spent on private tutoring in Europe each year"),
        ("4.7h", "avg. weekly study time saved with AI tools (pilot data)"),
    ]
    sw = Inches((12 - 0.4) / 3 - 0.2)
    sy = Inches(4.6)
    sh = Inches(2.0)
    for i, (big, small) in enumerate(stats):
        x = Inches(0.7) + i * (sw + Inches(0.2))
        add_rect(slide, x, sy, sw, sh, DEEP_RED)
        add_text(slide, x + Inches(0.2), sy + Inches(0.25),
                 sw - Inches(0.4), Inches(0.9), big,
                 size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.2), sy + Inches(1.2),
                 sw - Inches(0.4), Inches(0.7), small,
                 size=12, color=WHITE, align=PP_ALIGN.CENTER)


def slide_value(slide):
    add_text(slide, Inches(0.7), Inches(2.15), Inches(12), Inches(0.5),
             "Our customer will care about our solution because…",
             size=16, color=INK)

    # Filled mad-libs sentence
    add_text(slide, Inches(0.7), Inches(2.6), Inches(12), Inches(2.0),
             "It solves exam anxiety and confusion,\n"
             "and helps them learn faster and remember more\n"
             "without paying for tutors or wasting nights on YouTube.",
             size=22, bold=True, color=DEEP_RED)

    # Reasons-to-believe row
    reasons = [
        ("Personal", "Adapts to your course, your style, your language."),
        ("Affordable", "€7/month — less than one tutor hour."),
        ("Always on", "Answers at 2 a.m. the night before the exam."),
        ("Built-in", "Connects to Moodle, Canvas, and Google Drive."),
    ]
    cw = Inches((12 - 0.6) / 4 - 0.05)
    cy = Inches(5.3)
    ch = Inches(1.5)
    for i, (t, d) in enumerate(reasons):
        x = Inches(0.7) + i * (cw + Inches(0.2))
        add_rect(slide, x, cy, cw, ch, WHITE)
        add_text(slide, x + Inches(0.15), cy + Inches(0.15),
                 cw - Inches(0.3), Inches(0.4), t,
                 size=14, bold=True, color=DEEP_RED)
        add_text(slide, x + Inches(0.15), cy + Inches(0.55),
                 cw - Inches(0.3), Inches(0.9), d,
                 size=11, color=INK)


def slide_business(slide):
    body_x = Inches(0.7)
    # Left column - business model
    add_text(slide, body_x, Inches(2.15), Inches(6), Inches(0.5),
             "Business model", size=16, bold=True, color=DEEP_RED)
    add_bullets(slide, body_x, Inches(2.65), Inches(6), Inches(4),
                ["Freemium: free for 1 course, paid for unlimited.",
                 "Pro plan: €7/month for students.",
                 "Campus plan: bulk license sold to universities.",
                 "B2B add-on: white-label for online course providers.",
                 "Marketplace: verified tutors charge per session in-app."],
                size=14)

    # Right column - traction & next steps
    add_text(slide, Inches(7.0), Inches(2.15), Inches(5.7), Inches(0.5),
             "Traction & next 90 days", size=16, bold=True, color=DEEP_RED)
    add_bullets(slide, Inches(7.0), Inches(2.65), Inches(5.7), Inches(4),
                ["MVP live with 120 beta users from UE Berlin.",
                 "NPS 62 across first cohort.",
                 "Pilot agreement signed with 2 student associations.",
                 "Q3: launch on iOS + Android.",
                 "Q4: close pre-seed round (€350k)."],
                size=14)

    # Bottom CTA bar
    add_rect(slide, Inches(0.7), Inches(6.2), Inches(12), Inches(0.55), DEEP_RED)
    add_text(slide, Inches(0.7), Inches(6.27), Inches(12), Inches(0.45),
             "Looking for: pilot universities, design partners, and pre-seed angels.",
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def make_thanks_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, MINT)
    add_side_bar(slide)

    add_text(slide, Inches(0.7), Inches(1.8), Inches(12), Inches(1.6),
             "Thank you.", size=88, bold=True, color=DEEP_RED)
    add_text(slide, Inches(0.7), Inches(3.6), Inches(12), Inches(0.8),
             "Feedback directly afterwards.",
             size=28, bold=True, color=INK)
    add_text(slide, Inches(0.7), Inches(4.5), Inches(12), Inches(0.6),
             "Questions, critiques, brutal honesty — all welcome.",
             size=18, color=INK)

    # Contact strip
    add_rect(slide, Inches(0.7), Inches(5.8), Inches(12), Inches(1.0), WHITE)
    add_text(slide, Inches(0.9), Inches(5.92), Inches(11), Inches(0.4),
             "Team NOMIN  ·  hello@nomin.study",
             size=14, bold=True, color=DEEP_RED)
    add_text(slide, Inches(0.9), Inches(6.32), Inches(11), Inches(0.4),
             "University of Europe for Applied Sciences  ·  Entrepreneurship 2026",
             size=12, color=INK)

    add_footer(slide, 8, 8)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs)
    make_section_slide(prs, 1, "OUR START-UP IDEA IS…",
                       "NOMIN turns any course into a personal AI tutor.",
                       slide_idea, page=2)
    make_section_slide(prs, 2, "OUR PRIMARY CUSTOMER IS…",
                       "European university students who study smarter, not harder.",
                       slide_customer, page=3)
    make_section_slide(prs, 3, "THEIR MAIN NEED IS…",
                       "Understand the material, pass exams, stay sane.",
                       slide_need, page=4)
    make_section_slide(prs, 4, "OUR IMPACT",
                       "NOMIN reduces drop-outs and improves grades.",
                       slide_impact, page=5)
    make_section_slide(prs, 5, "WHY THEY'LL CARE",
                       "A personal tutor in your pocket — for the price of a coffee.",
                       slide_value, page=6)
    make_section_slide(prs, 6, "HOW WE MAKE IT WORK",
                       "Freemium for students, licenses for universities.",
                       slide_business, page=7)
    make_thanks_slide(prs)

    out = "/home/user/ilmkhona0.com.co/presentations/NOMIN_Startup_Pitch.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build()
